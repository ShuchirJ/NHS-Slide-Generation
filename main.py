from pptx import Presentation
import os

prs = Presentation("template.pptx")
slide_layout = prs.slide_layouts[0]

# placeholders:
# 0: title
# 10: picture

names = open("names.txt").read().replace(" , ", ", ").splitlines()
names = [name.split(", ") for name in names if name.strip() != ""]

photoGuide = open("11/11.TXT").read().splitlines()
photoGuide = [line.split("\t") for line in photoGuide if line.strip() != ""]

i = 1
for name in names:
    slide = prs.slides.add_slide(slide_layout)
    line = next((line for line in photoGuide if line[4] == name[0] and line[5] == name[1]), None)
    if line:
        filepath = "11/" + line[2]
    else:
        filepath = "11/" + input("Enter filename for " + name[1] + " " + name[0] + ": ")
    name = name[1] + " " + name[0]
    title = slide.shapes.title
    title.text = name

    picturePlaceholder = slide.placeholders[10]
    picturePlaceholder.insert_picture(filepath)
    
    print(f"Processed: {i}/{len(names)}")
    os.system('cls' if os.name == 'nt' else 'clear')
    i += 1

prs.save('out.pptx')